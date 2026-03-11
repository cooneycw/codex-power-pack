"""PowerPoint slide builder - creates PPTX files with embedded diagrams and content."""

from __future__ import annotations

import base64
import io
import logging
import re
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)


@dataclass
class SlideDefinition:
    """Typed slide definition with validation at construction time."""

    layout: str = "content"
    title: str = ""
    subtitle: str = ""
    content: str = ""
    left_content: str = ""
    right_content: str = ""
    image_base64: str = ""
    notes: str = ""

    _VALID_LAYOUTS: tuple[str, ...] = field(
        default=("title", "content", "diagram", "two-column", "blank"),
        init=False,
        repr=False,
    )

    def __post_init__(self) -> None:
        if self.layout not in self._VALID_LAYOUTS:
            logger.warning(
                f"Unknown layout '{self.layout}', defaulting to 'content'. "
                f"Valid layouts: {', '.join(self._VALID_LAYOUTS)}"
            )
            self.layout = "content"

    @classmethod
    def from_dict(cls, d: dict) -> SlideDefinition:
        """Create from a dict, ignoring unknown keys."""
        known = {f.name for f in cls.__dataclass_fields__.values() if f.init}
        return cls(**{k: v for k, v in d.items() if k in known})

# Maximum base64-encoded image size (10 MB decoded ~ 13.3 MB encoded)
MAX_IMAGE_BASE64_LENGTH = 14_000_000

# Terms that should not appear in generated output unless explicitly requested
_FRAMEWORK_NAMES = [
    "mckinsey", "bcg", "boston consulting", "bain", "deloitte", "accenture",
    "game framework", "pyramid principle", "minto pyramid",
    "scqa framework", "star framework", "monroe's motivated sequence",
]

_PLACEHOLDER_PATTERNS = [
    r"\[\s*insert\b[^\]]*\]",
    r"\[\s*placeholder\b[^\]]*\]",
    r"\[\s*todo\b[^\]]*\]",
    r"lorem ipsum",
]

# Standard 16:9 widescreen dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Blank layout index in the default python-pptx template
_BLANK_LAYOUT = 6


def _to_slide_def(item: dict | SlideDefinition) -> SlideDefinition:
    """Convert dict to SlideDefinition if needed."""
    if isinstance(item, SlideDefinition):
        return item
    return SlideDefinition.from_dict(item)


def create_presentation(
    title: str,
    slides: list[dict | SlideDefinition],
    author: str = "Nano-Banana",
) -> bytes:
    """Create a PowerPoint presentation from slide definitions.

    Args:
        title: Presentation title (used for title slide).
        slides: List of SlideDefinition or dicts with:
            - layout: "title", "content", "diagram", "two-column", "blank"
            - title: Slide title
            - subtitle: Optional subtitle
            - content: Text content (bullet points separated by newlines)
            - image_base64: Base64-encoded PNG image to embed
            - notes: Speaker notes
        author: Author metadata.

    Returns:
        PPTX file as bytes.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Set core properties
    prs.core_properties.author = author
    prs.core_properties.title = title

    for item in slides:
        sd = _to_slide_def(item)
        if sd.layout == "title":
            _add_title_slide(prs, sd)
        elif sd.layout == "diagram":
            _add_diagram_slide(prs, sd)
        elif sd.layout == "two-column":
            _add_two_column_slide(prs, sd)
        elif sd.layout == "blank":
            _add_blank_slide(prs, sd)
        else:
            _add_content_slide(prs, sd)

    # Save to bytes
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_title_slide(prs: Presentation, sd: SlideDefinition) -> None:
    """Add a title slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT])
    _set_slide_bg(slide, RGBColor(0x0F, 0x17, 0x2A))

    # Title
    _add_textbox(
        slide,
        sd.title,
        left=Inches(1), top=Inches(2.5), width=Inches(11.333), height=Inches(1.5),
        font_size=Pt(44), font_color=RGBColor(0xF8, 0xFA, 0xFC),
        bold=True, alignment=PP_ALIGN.CENTER,
    )

    # Subtitle
    if sd.subtitle:
        _add_textbox(
            slide,
            sd.subtitle,
            left=Inches(1), top=Inches(4.2), width=Inches(11.333), height=Inches(1),
            font_size=Pt(20), font_color=RGBColor(0x94, 0xA3, 0xB8),
            alignment=PP_ALIGN.CENTER,
        )

    _add_notes(slide, sd)


def _add_content_slide(prs: Presentation, sd: SlideDefinition) -> None:
    """Add a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT])
    _set_slide_bg(slide, RGBColor(0x1E, 0x29, 0x3B))

    # Title bar
    _add_textbox(
        slide,
        sd.title,
        left=Inches(0.8), top=Inches(0.5), width=Inches(11.733), height=Inches(0.9),
        font_size=Pt(32), font_color=RGBColor(0xF8, 0xFA, 0xFC),
        bold=True,
    )

    # Content
    if sd.content:
        _add_textbox(
            slide,
            sd.content,
            left=Inches(1), top=Inches(1.8), width=Inches(11.333), height=Inches(5),
            font_size=Pt(18), font_color=RGBColor(0xE2, 0xE8, 0xF0),
        )

    _add_notes(slide, sd)


def _add_diagram_slide(prs: Presentation, sd: SlideDefinition) -> None:
    """Add a slide with an embedded diagram image."""
    slide = prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT])
    _set_slide_bg(slide, RGBColor(0x0F, 0x17, 0x2A))

    # Title (smaller for diagram slides)
    if sd.title:
        _add_textbox(
            slide,
            sd.title,
            left=Inches(0.5), top=Inches(0.3), width=Inches(12.333), height=Inches(0.7),
            font_size=Pt(24), font_color=RGBColor(0xF8, 0xFA, 0xFC),
            bold=True,
        )

    # Embed image
    image_b64 = sd.image_base64
    if image_b64:
        if len(image_b64) > MAX_IMAGE_BASE64_LENGTH:
            error_msg = (
                f"Image too large ({len(image_b64):,} chars, "
                f"max {MAX_IMAGE_BASE64_LENGTH:,}). Reduce image size."
            )
            logger.warning(f"Skipping oversized image on diagram slide: {error_msg}")
            _add_textbox(
                slide,
                f"[{error_msg}]",
                left=Inches(1), top=Inches(3), width=Inches(11.333), height=Inches(1),
                font_size=Pt(16), font_color=RGBColor(0xEF, 0x44, 0x44),
            )
        else:
            try:
                img_bytes = base64.b64decode(image_b64)
                img_stream = io.BytesIO(img_bytes)
                # Full-width diagram below title
                slide.shapes.add_picture(
                    img_stream,
                    left=Inches(0.3),
                    top=Inches(1.2),
                    width=Inches(12.733),
                    height=Inches(5.8),
                )
            except (base64.binascii.Error, ValueError) as e:
                logger.warning(f"Failed to decode diagram image: {e}")
                _add_textbox(
                    slide,
                    f"[Diagram image could not be decoded: {e}]",
                    left=Inches(1), top=Inches(3), width=Inches(11.333), height=Inches(1),
                    font_size=Pt(16), font_color=RGBColor(0xEF, 0x44, 0x44),
                )
            except Exception as e:
                logger.warning(f"Failed to embed diagram image: {e}")
                _add_textbox(
                    slide,
                    f"[Diagram image could not be embedded: {e}]",
                    left=Inches(1), top=Inches(3), width=Inches(11.333), height=Inches(1),
                    font_size=Pt(16), font_color=RGBColor(0xEF, 0x44, 0x44),
                )

    _add_notes(slide, sd)


def _add_two_column_slide(prs: Presentation, sd: SlideDefinition) -> None:
    """Add a two-column layout slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT])
    _set_slide_bg(slide, RGBColor(0x1E, 0x29, 0x3B))

    # Title
    _add_textbox(
        slide,
        sd.title,
        left=Inches(0.8), top=Inches(0.5), width=Inches(11.733), height=Inches(0.9),
        font_size=Pt(32), font_color=RGBColor(0xF8, 0xFA, 0xFC),
        bold=True,
    )

    # Left column
    left_content = sd.left_content or sd.content
    _add_textbox(
        slide,
        left_content,
        left=Inches(0.8), top=Inches(1.8), width=Inches(5.5), height=Inches(5),
        font_size=Pt(16), font_color=RGBColor(0xE2, 0xE8, 0xF0),
    )

    # Right column
    if sd.right_content:
        _add_textbox(
            slide,
            sd.right_content,
            left=Inches(7), top=Inches(1.8), width=Inches(5.5), height=Inches(5),
            font_size=Pt(16), font_color=RGBColor(0xE2, 0xE8, 0xF0),
        )

    _add_notes(slide, sd)


def _add_blank_slide(prs: Presentation, sd: SlideDefinition) -> None:
    """Add a blank slide (dark background only)."""
    slide = prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT])
    _set_slide_bg(slide, RGBColor(0x0F, 0x17, 0x2A))
    _add_notes(slide, sd)


def _set_slide_bg(slide, color: RGBColor) -> None:
    """Set slide background to a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    font_size=Pt(18),
    font_color=RGBColor(0xE2, 0xE8, 0xF0),
    bold: bool = False,
    alignment=PP_ALIGN.LEFT,
) -> None:
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line.lstrip("- ")
        p.alignment = alignment
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.font.bold = bold

        # Add bullet for lines starting with -
        if line.strip().startswith("-"):
            p.level = 0


def _add_notes(slide, sd: SlideDefinition) -> None:
    """Add speaker notes to a slide."""
    if sd.notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = sd.notes


def validate_slides(
    slides: list[dict | SlideDefinition],
    prohibited_terms: list[str] | None = None,
) -> dict:
    """Validate slide definitions before PPTX creation.

    Checks for framework name leaks, placeholder text, structural issues,
    and content quality. Returns pass/fail with specific issues.

    Args:
        slides: List of SlideDefinition or dicts (same format as create_presentation).
        prohibited_terms: Additional terms to flag (case-insensitive).

    Returns:
        dict with passed (bool), issues (list), and summary.
    """
    issues: list[dict] = []
    terms = _FRAMEWORK_NAMES + (prohibited_terms or [])

    _TEXT_FIELDS = ("title", "subtitle", "content", "left_content", "right_content", "notes")

    for i, item in enumerate(slides, 1):
        sd = _to_slide_def(item)

        # Collect all text from this slide
        text_fields = [
            (key, getattr(sd, key))
            for key in _TEXT_FIELDS
            if getattr(sd, key)
        ]

        # Check: framework/corporate name leaks (word-boundary match)
        for field_name, text in text_fields:
            text_lower = text.lower()
            for term in terms:
                pattern = r"\b" + re.escape(term.lower()) + r"\b"
                if re.search(pattern, text_lower):
                    issues.append({
                        "slide": i,
                        "field": field_name,
                        "severity": "high",
                        "check": "framework_attribution",
                        "message": f"Found '{term}' in {field_name}. Remove framework/corporate name attribution.",
                    })

        # Check: placeholder text
        for field_name, text in text_fields:
            for pat in _PLACEHOLDER_PATTERNS:
                matches = re.findall(pat, text, re.IGNORECASE)
                for match in matches:
                    issues.append({
                        "slide": i,
                        "field": field_name,
                        "severity": "high",
                        "check": "placeholder_text",
                        "message": f"Found placeholder text '{match}' in {field_name}.",
                    })

        # Check: diagram slide must have image
        if sd.layout == "diagram" and not sd.image_base64:
            issues.append({
                "slide": i,
                "field": "image_base64",
                "severity": "medium",
                "check": "missing_image",
                "message": "Diagram slide has no image_base64. Will render as empty slide.",
            })

        # Check: content slide should have content
        if sd.layout == "content" and not sd.content:
            issues.append({
                "slide": i,
                "field": "content",
                "severity": "low",
                "check": "empty_content",
                "message": "Content slide has no body text.",
            })

        # Check: two-column should have both columns
        if sd.layout == "two-column":
            if not sd.left_content and not sd.content:
                issues.append({
                    "slide": i,
                    "field": "left_content",
                    "severity": "low",
                    "check": "empty_column",
                    "message": "Two-column slide missing left column content.",
                })
            if not sd.right_content:
                issues.append({
                    "slide": i,
                    "field": "right_content",
                    "severity": "low",
                    "check": "empty_column",
                    "message": "Two-column slide missing right column content.",
                })

    high_count = sum(1 for iss in issues if iss["severity"] == "high")
    passed = high_count == 0

    return {
        "passed": passed,
        "slide_count": len(slides),
        "issue_count": len(issues),
        "high_severity": high_count,
        "issues": issues,
        "summary": "All checks passed." if passed else f"{high_count} high-severity issue(s) must be fixed.",
    }
